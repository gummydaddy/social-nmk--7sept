# from datetime import timedelta, timezone
from urllib import request
import uuid
import zipfile
import pyrebase
from django import forms
from django.forms import CharField, EmailField
from django.shortcuts import render, redirect, get_object_or_404
from django.contrib.auth import authenticate, login, logout
from django.contrib.auth.models import User  # Import the User model at the top
from django.http import HttpResponse, HttpResponseBadRequest, JsonResponse, StreamingHttpResponse, HttpResponseNotFound, HttpResponseServerError
from django.contrib.auth.decorators import login_required, user_passes_test
from django.contrib.admin.views.decorators import staff_member_required
from django.contrib import messages
from django.core.mail import send_mail, EmailMessage, get_connection
from django.core.cache import cache
from django.utils.module_loading import import_string
from django.conf import settings
from .models import KYC, File, TemporaryUser, UserAssociation, UserUpload, RegistrationForm, CustomGroup, CustomGroupAdmin#, TemporarilyLock
from .forms import CustomSignupForm, UserUploadForm, DeleteUploadForm, RegistrationFormForm, KYCForm, CardForm, GroupCreationForm, SubgroupSignupForm, PasswordResetForm
from twilio.rest import Client
from cryptography.fernet import Fernet
import os
import logging
from django.utils import timezone  # Import Django's timezone utility
from datetime import timedelta
from django.contrib.auth import update_session_auth_hash
from django.utils.http import urlsafe_base64_encode, urlsafe_base64_decode
import mimetypes
from docx import Document
from pptx import Presentation
import xml.etree.ElementTree as ET
from django.template.response import TemplateResponse
from django.core.files.base import ContentFile
import openpyxl  # To handle .xlsx files
import xlrd  # To handle .xls files
# import base64
from PyPDF2 import PdfReader
from django.core.files.storage import FileSystemStorage
from django.utils.html import escape
from pdf2image import convert_from_path
from PIL import Image
import io
import shutil
from .tasks import process_file_upload  # import your Celery task
from django.views.decorators.cache import never_cache









firebaseConfig={
    'apiKey': "AIzaSyCN7s6q2iX84wAa-bcADtojDbRKRsK5UTk",
    'authDomain': "authporp.firebaseapp.com",
    'projectId': "authporp",
    'storageBucket': "authporp.appspot.com",
    'messagingSenderId': "746007263598",
    'appId': "1:746007263598:web:1d9cbb363e41cba3d3e97f",
    'measurementId': "G-MQ4Q1F4FCE",
    "databaseURL": "",  # Ensure this is included
}

firebase= pyrebase.initialize_app(firebaseConfig)
auth=firebase.auth()


def home(request):
    return render(request, 'home.html')


def TermAndCondition(request):
    return render(request, 'TermAndCondition.html')


def send_confirmation_email(user):
    subject = "Welcome to NMK Financial Services - Django Login!"
    message = f"Hello {user.first_name},\nWelcome to NMK!\nThank you for being a part of our community."
    email_from = settings.EMAIL_HOST_USER
    recipient_list = [user.email]
    send_mail(subject, message, email_from, recipient_list)




def signup(request):
    if request.method == 'POST':
        form = CustomSignupForm(request.POST)
        if form.is_valid():
            # Extract user data from the form
            username = form.cleaned_data.get('username')
            email = form.cleaned_data.get('email')
            password = form.cleaned_data.get('password1')  # Assuming you have a password1 field
            first_name = form.cleaned_data.get('first_name')
            last_name = form.cleaned_data.get('last_name')

            try:
                # Create user with Firebase
                firebase_user = auth.create_user_with_email_and_password(email, password)
                # Send email verification
                auth.send_email_verification(firebase_user['idToken'])

                # Save user details in Django model without committing to the database
                user_model = form.save(request, commit=False)  # Pass request to save method
                user_model.first_name = first_name
                user_model.last_name = last_name
                user_model.email = email
                user_model.username = username  # Set username to username 
                user_model.save()  # Now commit to the database

                # Log the user in Django
                login(request, user_model)
                messages.success(request, "Your account has been created successfully. Please check your email to confirm your account.")
                return redirect('/')
            except Exception as e:
                # Handle Firebase authentication errors
                messages.error(request, f"Failed to create account: {e}")
        else:
            messages.error(request, "Failed to create account. Please check the form entries.")
    else:
        form = CustomSignupForm()
    return render(request, 'signup.html', {'form': form})


@never_cache
def login_view(request):
    if request.user.is_authenticated:

        return redirect('/following_media')  

    if request.method == 'POST':
        username_or_email = request.POST['username']
        password = request.POST['pass1']
        remember_me = request.POST.get('remember_me')  # This is a checkbox in your HTML

        try:
            # Check if the input is an email or username
            if '@' in username_or_email:
                # Authenticate with email
                user = User.objects.get(email=username_or_email)
                username = user.username
            else:
                # Authenticate with username
                username = username_or_email
                user = User.objects.get(username=username)

            # Firebase authentication: Check if the provided credentials are correct in Firebase
            try:
                firebase_user = auth.sign_in_with_email_and_password(user.email, password)

                # At this point, the password matches Firebase. Update Django's password with Firebase password.
                user.set_password(password)  # Sync the Firebase password with Django's password
                user.save()

                # Authenticate the user in Django with the updated password
                user = authenticate(request, username=username, password=password)

                if user is not None:
                    # Log the user in Django
                    login(request, user)

                    # Ensure the session remains intact after password update
                    update_session_auth_hash(request, user)
                    # Set session expiry based on "remember me"
                    if not remember_me:
                        request.session.set_expiry(0)  # Session expires on browser/app close
                    else:
                        request.session.set_expiry(60 * 60 * 24 * 14)  # 14 days

                    response = redirect('/following_media')  # Default redirect

                    # Set session cookie and cache the username
                    #response = HttpResponse("You're logged in.")
                    #response.set_cookie('username', username)
                    #cache.set(f'user_{user.id}', user.username)

                    # Redirect based on user roles
                    if CustomGroupAdmin.objects.filter(user=user).exists():
                        return redirect('/subgroup_landing_page')
                    elif user.is_staff:
                        return redirect('/super_user_landing_page')
                    #else:
                        # return redirect('/landing_page')
                        #return redirect('/following_media')
                    # Set session/cookie/cache
                    response.set_cookie('username', username)
                    cache.set(f'user_{user.id}', username)
                    return response  

            except Exception as firebase_error:
                # Handle Firebase authentication errors
                error_message = str(firebase_error)
                if "EMAIL_NOT_FOUND" in error_message:
                    messages.error(request, "No account found with this email.")
                elif "INVALID_PASSWORD" in error_message:
                    messages.error(request, "The password is incorrect. Please try again.")
                else:
                    messages.error(request, f"Firebase error: {error_message}")

        except User.DoesNotExist:
            messages.error(request, 'Invalid username or email')

        #except Exception as e:
            # Handle any other errors
            #messages.error(request, f"Failed to log in: {str(e)}")

        # Group-based authentication (if necessary)
        try:
            group = CustomGroup.objects.get(name=username_or_email)
            user = group.users.first()
            if user:
                login(request, user)
                # Set session expiry based on "remember me" here as well
                if not remember_me:
                    request.session.set_expiry(0)
                else:
                    request.session.set_expiry(60 * 60 * 24 * 14)

                request.session['association_name'] = username_or_email
                request.session['user_id'] = user.id
                cache.set(f'user_{user.id}', username_or_email)
                return redirect('/landing_page')
        except CustomGroup.DoesNotExist:
            #pass
            messages.error(request, 'Invalid username, email, or group name.')

        except Exception as e:
            messages.error(request, f"Login failed: {str(e)}")

        return render(request, 'login_view.html')

    else:
        return render(request, 'login_view.html')


# Cache settings
# CACHE_TIMEOUT = 300  # 5 minutes
# CACHE_PREFIX = 'auth'

# def get_cache_key(key):
#     return f"{CACHE_PREFIX}_{key}"

# def authenticate_user(username, password):
#     """
#     Authenticate user using Firebase and Django.
    
#     Args:
#     username (str): User's username or email.
#     password (str): User's password.
    
#     Returns:
#     User or None: Authenticated user object or None.
#     """
#     cache_key = 'my_key'
#     user = cache.get(cache_key)
    
#     if not user:
#         try:
#             # Check if the input is an email or username
#             if '@' in username:
#                 user = User.objects.get(email=username)
#             else:
#                 user = User.objects.get(username=username)
            
#             # Cache user object
#             cache.set(cache_key, user, CACHE_TIMEOUT)
#         except User.DoesNotExist:
#             return None
        
#         # Firebase authentication
#         try:
#             firebase_user = auth.sign_in_with_email_and_password(user.email, password)
#             user.set_password(password)
#             user.save()
#         except Exception as e:
#             # Handle Firebase authentication errors
#             return None
    
#     return user

# def group_based_authentication(username_or_email):
#     """
#     Authenticate user using group-based authentication.
    
#     Args:
#     username_or_email (str): Group name or username.
    
#     Returns:
#     User or None: Authenticated user object or None.
#     """
#     cache_key = get_cache_key(f"group_{username_or_email}")
#     user = cache.get(cache_key)
    
#     if not user:
#         try:
#             group = CustomGroup.objects.get(name=username_or_email)
#             user = group.users.first()
            
#             # Cache user object
#             cache.set(cache_key, user, CACHE_TIMEOUT)
#         except CustomGroup.DoesNotExist:
#             return None
    
#     return user

# def login_view(request):
#     if request.method == 'POST':
#         username_or_email = request.POST['username']
#         password = request.POST['pass1']
        
#         user = authenticate_user(username_or_email, password)
#         if user is None:
#             user = group_based_authentication(username_or_email)
        
#         if user:
#             login(request, user)
#             # Update session and cache
#             cache.set(get_cache_key(f"user_{user.id}"), user.username, CACHE_TIMEOUT)
#             return redirect('/landing_page')
#         else:
#             # Handle authentication failure
#             messages.error(request, 'Invalid username or password')
#             return render(request, 'login_view.html')
#     else:
#         return render(request, 'login_view.html')



@staff_member_required
def super_user_landing_page(request):
    registration_forms = RegistrationForm.objects.all()
    return render(request, 'super_user_landing_page.html', {'registration_forms': registration_forms})

@login_required
def subgroup_landing_page(request):
    user_admin_groups = request.user.admin_groups.filter(is_approved=True)
    approved_groups = CustomGroup.objects.filter(is_approved=True)
    pending_requests = RegistrationForm.objects.filter(status='pending')
    context = {
        'user_admin_groups': user_admin_groups,
        'approved_groups': approved_groups,
        'pending_requests': pending_requests,
    }
    return render(request, 'subgroup_landing_page.html', context)


# LogOut
@login_required
def logout_view(request):
    user = request.user
    if user.is_authenticated:
        cache.delete(f'user_{user.id}')
    request.session.flush()
    logout(request)
   # request.session.flush()
    return redirect('/')


def password_reset(request):
    if request.method == 'POST':
        form = PasswordResetForm(request.POST)
        if form.is_valid():
            username_or_email = form.cleaned_data['username_or_email']
            try:
                # Check if username_or_email is a username or email
                if '@' in username_or_email:
                    # It's an email
                    user = User.objects.get(email=username_or_email)
                else:
                    # It's a username
                    user = User.objects.get(username=username_or_email)
                # Send Firebase password reset email
                auth.send_password_reset_email(user.email)
                messages.success(request, 'A password reset link has been sent to your email.')
                return redirect('/login/')
            except User.DoesNotExist:
                messages.error(request, "No account found with this username or email address.")
            except Exception as e:
                error_message = str(e)
                messages.error(request, f"Error: {error_message}")
        else:
            messages.error(request, 'Please provide a valid username or email.')
    else:
        form = PasswordResetForm()
        if request.user.is_authenticated:
            email = request.user.email
            form.fields['username_or_email'].initial = email
            form.fields['username_or_email'].widget = forms.HiddenInput()

    return render(request, 'password_reset.html', {'form': form})



# Landing
@login_required   
def landing_page(request):
    cache_key = f'user_{request.user.id}_username'
    user_username = cache.get(cache_key)
    #user_username = cache.get(f'user_{request.user.id}')
    if not user_username:
        user_username = request.user.username
        cache.set(cache_key, user_username, timeout=60 * 60 * 24)  # Cache for 1 day
        #cache.set(f'user_{request.user.id}', user_username, timeout=3600)
    try:
        user_card = request.user.card
    except User.card.RelatedObjectDoesNotExist:
        user_card = None
    user_uploads = UserUpload.objects.filter(user=request.user).order_by('-upload_date')[:]
    super_user = request.session.get('super_user', False)
    context = {
        'user_username': user_username,
        'user_uploads': user_uploads,
        'super_user': super_user,
        'user_card': user_card,
        'user_id': request.user.id,
    }
    return render(request, 'landing_page.html', context)

@login_required
def group_list(request):
    groups = CustomGroup.objects.all()
    return render(request, 'group_list.html', {'groups': groups})

@login_required
def group_create(request):
    if request.method == 'POST':
        form = GroupCreationForm(request.POST)
        if form.is_valid():
            group = form.save()
            return redirect('/group_list',{'group':group})
    else:
        form = GroupCreationForm()
    return render(request, 'customgroup_form.html', {'form': form})

@login_required
def subgroup_signup(request):
    if request.method == 'POST':
        form = SubgroupSignupForm(request.POST, request.FILES, user=request.user)
        if form.is_valid():
            form.save()
            messages.success(request, 'Your subgroup signup request has been submitted for approval.')
            return redirect('only_card:landing_page')
    else:
        form = SubgroupSignupForm(user=request.user)
    return render(request, 'subgroup_signup_form.html', {'form': form})

'''
@login_required
def upload_document(request):
    if request.method == 'POST':
        form = UserUploadForm(request.POST, request.FILES)
        if form.is_valid():
            form.instance.user = request.user

            try:
                form.save()  # This will check the storage limit
            except ValueError as e:
                # If storage is full, show a message and prevent upload
                messages.error(request, "Storage full. Please buy more storage.")
                return redirect('/buy_storage')  # Redirect to a "Buy Storage" page

            if CustomGroupAdmin.objects.filter(user=request.user).exists():
                return redirect('/subgroup_landing_page')
            else:
                return redirect('/landing_page')
    else:
        form = UserUploadForm()

    return render(request, 'upload_document.html', {'form': form})
'''

def buy_storage(request):
    return render(request, 'buy_storage.html')


def upload_document(request):
    if request.method == 'POST':
        form = UserUploadForm(request.POST, request.FILES)
        if form.is_valid():
            upload_instance = form.save(commit=False)
            upload_instance.user = request.user
            upload_instance.save()  # Save metadata without full validation

            # Call the Celery task
            process_file_upload.delay(upload_instance.id)

            messages.success(request, "Upload started! We'll notify you once it's processed.")

            if CustomGroupAdmin.objects.filter(user=request.user).exists():
                return redirect('/subgroup_landing_page')
            else:
                return redirect('/landing_page')
    else:
        form = UserUploadForm()

    return render(request, 'upload_document.html', {'form': form})



# Set up logging
logger = logging.getLogger(__name__)

@login_required
def view_file(request, upload_id):
    try:
        # Fetch the uploaded file record for the authenticated user
        upload = get_object_or_404(UserUpload, id=upload_id, user=request.user)

        # Ensure the file has an encryption key
        if not upload.encryption_key:
            logger.error(f"Encryption key not found for file with ID {upload_id}")
            return HttpResponseServerError("Encryption key not found for this file.")

        try:
            # Initialize the Fernet cipher for decryption using the encryption key
            fernet = Fernet(upload.encryption_key.encode('utf-8'))
        except Exception as e:
            logger.error(f"Error initializing Fernet cipher: {e}")
            return HttpResponseServerError("Failed to initialize decryption.")

        # Check if the file exists on the server
        file_path = upload.file.path
        if not os.path.exists(file_path):
            logger.error(f"File not found: {file_path}")
            return HttpResponseNotFound("File not found")
        

        # Detect the correct MIME type based on the file extension
        mime_type, _ = mimetypes.guess_type(file_path)
        if not mime_type:
            mime_type = 'application/octet-stream'  # Fallback to a default type

        try:
            # Read the entire encrypted file into memory
            with open(file_path, 'rb') as encrypted_file:
                encrypted_file_data = encrypted_file.read()

            # Decrypt the entire file at once
            decrypted_file_data = fernet.decrypt(encrypted_file_data)

            # Return the decrypted file as a streaming HTTP response
            response = StreamingHttpResponse(
                iter([decrypted_file_data]),  # Send the decrypted content as a single iterable
                content_type=mime_type
            )
            response['Content-Disposition'] = f'attachment; filename="{upload.file_name}"'
            return response

        except Exception as e:
            logger.error(f"Error during file decryption: {e}")
            return HttpResponseServerError("Error decrypting file.")
            

    except Exception as e:
        logger.error(f"Unhandled exception in view_file: {e}")
        return HttpResponseServerError("An error occurred while processing your request.")



@login_required
def view_pdf_file(request, upload_id):
    try:
        # Fetch the uploaded file record for the authenticated user
        upload = get_object_or_404(UserUpload, id=upload_id, user=request.user)

        # Ensure the file has an encryption key
        if not upload.encryption_key:
            logger.error(f"Encryption key not found for file with ID {upload_id}")
            return HttpResponseServerError("Encryption key not found for this file.")

        try:
            # Initialize the Fernet cipher for decryption using the encryption key
            fernet = Fernet(upload.encryption_key.encode('utf-8'))
        except Exception as e:
            logger.error(f"Error initializing Fernet cipher: {e}")
            return HttpResponseServerError("Failed to initialize decryption.")

        # Check if the file exists on the server
        file_path = upload.file.path
        if not os.path.exists(file_path):
            logger.error(f"File not found: {file_path}")
            return HttpResponseNotFound("File not found")

        # Check if the file is a .pdf file
        if not file_path.endswith('.pdf'):
            logger.error(f"Unsupported file type: {file_path}")
            return HttpResponseServerError("Unsupported file type.")

        # Read and decrypt the entire .pdf file
        try:
            with open(file_path, 'rb') as encrypted_file:
                encrypted_file_data = encrypted_file.read()

            decrypted_file_data = fernet.decrypt(encrypted_file_data)

            # Save the decrypted data to a temporary file to be read by PyPDF2
            temp_file_path = os.path.join('/tmp', f"decrypted_{upload.file_name}")
            with open(temp_file_path, 'wb') as temp_file:
                temp_file.write(decrypted_file_data)

            fs = FileSystemStorage('/tmp')
            with fs.open(f"decrypted_{upload.file_name}", 'rb') as pdf:
                response = HttpResponse(pdf, content_type='application/pdf')
                response['Content-Disposition'] = f'inline; filename="{upload.file_name}"'

            # Clean up: remove the temporary file after processing
            os.remove(temp_file_path)

            return response

        except Exception as e:
            logger.error(f"Error reading .pdf file: {e}")
            return HttpResponseServerError("Error reading .pdf file.")

    except Exception as e:
        logger.error(f"Unhandled exception in view_pdf_file: {e}")
        return HttpResponseServerError("An error occurred while processing your request.")


@login_required
def view_docx_file(request, upload_id):
    try:
        # Fetch the uploaded file record for the authenticated user
        upload = get_object_or_404(UserUpload, id=upload_id, user=request.user)

        # Ensure the file has an encryption key
        if not upload.encryption_key:
            logger.error(f"Encryption key not found for file with ID {upload_id}")
            return HttpResponseServerError("Encryption key not found for this file.")

        try:
            # Initialize the Fernet cipher for decryption using the encryption key
            fernet = Fernet(upload.encryption_key.encode('utf-8'))
        except Exception as e:
            logger.error(f"Error initializing Fernet cipher: {e}")
            return HttpResponseServerError("Failed to initialize decryption.")

        # Check if the file exists on the server
        file_path = upload.file.path
        if not os.path.exists(file_path):
            logger.error(f"File not found: {file_path}")
            return HttpResponseNotFound("File not found")

        # Check if the file is a .docx file
        if not file_path.endswith('.docx'):
            logger.error(f"Unsupported file type: {file_path}")
            return HttpResponseServerError("Unsupported file type.")

        # Read and decrypt the entire .docx file
        try:
            with open(file_path, 'rb') as encrypted_file:
                encrypted_file_data = encrypted_file.read()

            decrypted_file_data = fernet.decrypt(encrypted_file_data)

            # Save the decrypted data to a temporary file to be read by python-docx
            temp_file_path = os.path.join('/tmp', f"decrypted_{upload.file_name}")
            with open(temp_file_path, 'wb') as temp_file:
                temp_file.write(decrypted_file_data)

            # Use python-docx to open and read the content of the .docx file
            document = Document(temp_file_path)

            # Extract text from the .docx file and convert it to HTML format
            docx_content = ""
            for paragraph in document.paragraphs:
                docx_content += f"<p>{paragraph.text}</p>"

            # Clean up: remove the temporary file after processing
            os.remove(temp_file_path)

            # Render the content in an HTML template
            return render(request, 'view_docx.html', {'docx_content': docx_content, 'file_name': upload.file_name})

        except Exception as e:
            logger.error(f"Error reading .docx file: {e}")
            return HttpResponseServerError("Error reading .docx file.")

    except Exception as e:
        logger.error(f"Unhandled exception in view_docx_file: {e}")
        return HttpResponseServerError("An error occurred while processing your request.")


@login_required
def view_pptx_file(request, upload_id):
    try:
        # Fetch the uploaded file record for the authenticated user
        upload = get_object_or_404(UserUpload, id=upload_id, user=request.user)

        # Ensure the file has an encryption key
        if not upload.encryption_key:
            logger.error(f"Encryption key not found for file with ID {upload_id}")
            return HttpResponseServerError("Encryption key not found for this file.")

        try:
            # Initialize the Fernet cipher for decryption using the encryption key
            fernet = Fernet(upload.encryption_key.encode('utf-8'))
        except Exception as e:
            logger.error(f"Error initializing Fernet cipher: {e}")
            return HttpResponseServerError("Failed to initialize decryption.")

        # Check if the file exists on the server
        file_path = upload.file.path
        if not os.path.exists(file_path):
            logger.error(f"File not found: {file_path}")
            return HttpResponseNotFound("File not found")

        # Check if the file is a .pptx file
        if not file_path.endswith('.pptx'):
            logger.error(f"Unsupported file type: {file_path}")
            return HttpResponseServerError("Unsupported file type.")

        # Read and decrypt the entire .pptx file
        try:
            with open(file_path, 'rb') as encrypted_file:
                encrypted_file_data = encrypted_file.read()

            decrypted_file_data = fernet.decrypt(encrypted_file_data)

            # Save the decrypted data to a temporary file to be read by python-pptx
            temp_file_path = os.path.join('/tmp', f"decrypted_{upload.file_name}")
            with open(temp_file_path, 'wb') as temp_file:
                temp_file.write(decrypted_file_data)

            # Use python-pptx to open and read the content of the .pptx file
            presentation = Presentation(temp_file_path)

            # Extract slide content
            pptx_content = ""
            for slide_number, slide in enumerate(presentation.slides, start=1):
                pptx_content += f"<h3>Slide {slide_number}</h3>"
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        pptx_content += f"<p>{shape.text}</p>"

            # Clean up: remove the temporary file after processing
            os.remove(temp_file_path)

            # Render the content in an HTML template
            return render(request, 'view_docx.html', {'pptx_content': pptx_content, 'file_name': upload.file_name})

        except Exception as e:
            logger.error(f"Error reading .pptx file: {e}")
            return HttpResponseServerError("Error reading .pptx file.")

    except Exception as e:
        logger.error(f"Unhandled exception in view_pptx_file: {e}")
        return HttpResponseServerError("An error occurred while processing your request.")


def is_valid_excel_file(file_path):
    # Supported MIME types for Excel files
    excel_mime_types = [
        'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet',  # .xlsx
        'application/vnd.ms-excel.sheet.macroEnabled.12',  # .xlsm
        'application/vnd.openxmlformats-officedocument.spreadsheetml.template',  # .xltx
        'application/vnd.ms-excel.template.macroEnabled.12'  # .xltm
    ]
    
    mime_type, _ = mimetypes.guess_type(file_path)
    return mime_type in excel_mime_types



def view_xml_file(request, upload_id):
    try:
        # Fetch the uploaded file record for the authenticated user
        upload = get_object_or_404(UserUpload, id=upload_id, user=request.user)

        # Ensure the file has an encryption key
        if not upload.encryption_key:
            logger.error(f"Encryption key not found for file with ID {upload_id}")
            return HttpResponseServerError("Encryption key not found for this file.")

        try:
            # Initialize the Fernet cipher for decryption using the encryption key
            fernet = Fernet(upload.encryption_key.encode('utf-8'))
        except Exception as e:
            logger.error(f"Error initializing Fernet cipher: {e}")
            return HttpResponseServerError("Failed to initialize decryption.")

        # Check if the file exists on the server
        file_path = upload.file.path
        if not os.path.exists(file_path):
            logger.error(f"File not found: {file_path}")
            return HttpResponseNotFound("File not found")

        # Get the file extension and handle accordingly
        file_extension = os.path.splitext(file_path)[1].lower()

        # Decrypt and handle the file based on its extension
        if file_extension == '.xml':
            return handle_xml_or_xlsx_file(file_path, fernet, upload, file_type='xml')
        elif file_extension == '.xlsx' or file_extension == '.xls':
            return handle_xml_or_xlsx_file(file_path, fernet, upload, file_type=file_extension)
        else:
            logger.error(f"Unsupported file type: {file_path}")
            return HttpResponseServerError("Unsupported file type.")

    except Exception as e:
        logger.error(f"Unhandled exception in view_xml_file: {e}")
        return HttpResponseServerError("An error occurred while processing your request.")


def handle_xml_or_xlsx_file(file_path, fernet, upload, file_type):
    """Handles XML, XLS, and XLSX file decryption and rendering."""
    try:
        # Read and decrypt the entire file
        with open(file_path, 'rb') as encrypted_file:
            encrypted_file_data = encrypted_file.read()

        # Write the decrypted content to a temporary file for processing
        decrypted_file_path = os.path.join('/tmp', f"decrypted_{upload.file_name}")
        with open(decrypted_file_path, 'wb') as decrypted_file:
            decrypted_file.write(fernet.decrypt(encrypted_file_data))

        # Process based on file type
        if file_type == 'xml':
            # Handle XML file
            return handle_xml_file(decrypted_file_path, upload)

        elif file_type == '.xlsx':
            # Handle XLSX file
            return handle_xlsx_file(decrypted_file_path, upload)

        elif file_type == '.xls':
            # Handle XLS file using xlrd
            return handle_xls_file(decrypted_file_path, upload)

        else:
            logger.error("Unsupported file type provided.")
            return HttpResponseServerError("Unsupported file type.")

    except Exception as e:
        logger.error(f"Error decrypting {file_type} file: {e}")
        return HttpResponseServerError(f"Error decrypting {file_type} file.")
    finally:
        # Clean up the temporary decrypted file if it exists
        if os.path.exists(decrypted_file_path):
            os.remove(decrypted_file_path)


def handle_xml_file(decrypted_file_path, upload):
    """Handles the XML file content."""
    try:
        # Parse the XML content
        with open(decrypted_file_path, 'r', encoding='utf-8') as decrypted_file:
            xml_data = decrypted_file.read()

        root = ET.fromstring(xml_data)

        # Convert XML content to a string for display
        xml_content = ET.tostring(root, encoding='unicode', method='xml')

        # Render the content in an HTML template
        return render(request, 'view_docx.html', {
            'xml_content': xml_content,
            'file_name': upload.file_name
        })

    except ET.ParseError as e:
        logger.error(f"Error parsing XML file: {e}")
        return HttpResponseServerError("Error parsing the XML file.")


def handle_xlsx_file(decrypted_file_path, upload):
    """Handles the XLSX file content using openpyxl."""
    try:
        # Log the decrypted file path to ensure decryption is happening
        logger.info(f"Decrypted file path: {decrypted_file_path}")
        
        # Check the file size to make sure it's a valid Excel file size
        file_size = os.path.getsize(decrypted_file_path)
        logger.info(f"Decrypted file size: {file_size} bytes")
        
        # Open the decrypted Excel file using openpyxl
        workbook = openpyxl.load_workbook(decrypted_file_path)
        sheet = workbook.active

        # Collect data from the first sheet
        excel_data = []
        for row in sheet.iter_rows(values_only=True):
            excel_data.append(row)

        # Render the content in an HTML template
        return render(request, 'view_docx.html', {
            'excel_data': excel_data,
            'file_name': upload.file_name
        })

    except Exception as e:
        logger.error(f"Error reading Excel file (.xlsx): {e}")
        return HttpResponseServerError("Error reading Excel file (.xlsx).")



def handle_xls_file(decrypted_file_path, upload):
    """Handles the older .xls file format using xlrd."""
    try:
        # Open the decrypted Excel file using xlrd
        workbook = xlrd.open_workbook(decrypted_file_path)
        sheet = workbook.sheet_by_index(0)

        # Collect data from the first sheet
        excel_data = []
        for row_num in range(sheet.nrows):
            row = sheet.row_values(row_num)
            excel_data.append(row)

        # Render the content in an HTML template
        return render(request, 'view_docx.html', {
            'excel_data': excel_data,
            'file_name': upload.file_name
        })

    except Exception as e:
        logger.error(f"Error reading Excel file (.xls): {e}")
        return HttpResponseServerError("Error reading Excel file (.xls).")
    

@login_required
def view_text_file(request, upload_id):
    try:
        # Fetch the uploaded file record for the authenticated user
        upload = get_object_or_404(UserUpload, id=upload_id, user=request.user)

        # Ensure the file has an encryption key
        if not upload.encryption_key: 
            logger.error(f"Encryption key not found for file with ID {upload_id}")
            return HttpResponseServerError("Encryption key not found for this file.")

        try:
            # Initialize the Fernet cipher for decryption using the encryption key
            fernet = Fernet(upload.encryption_key.encode('utf-8'))
        except Exception as e:
            logger.error(f"Error initializing Fernet cipher: {e}")
            return HttpResponseServerError("Failed to initialize decryption.")

        # Check if the file exists on the server
        file_path = upload.file.path
        if not os.path.exists(file_path):
            logger.error(f"File not found: {file_path}")
            return HttpResponseNotFound("File not found")

        # Read and decrypt the entire file
        try:
            with open(file_path, 'rb') as encrypted_file:
                encrypted_file_data = encrypted_file.read()

            decrypted_file_data = fernet.decrypt(encrypted_file_data).decode('utf-8')

            # Determine the file type and render the appropriate template
            file_extension = os.path.splitext(upload.file.name)[1].lower()
            
            if file_extension == '.py':
                file_type = 'python'
            elif file_extension == '.js':
                file_type = 'javascript'
            elif file_extension == '.txt':
                file_type = 'text'
            else:
                file_type = 'unknown'

            # Render the content in an HTML template, pass decrypted content
            return render(request, 'view_docx.html', {
                'txt_content': decrypted_file_data,
                'file_name': upload.file_name,
                'file_type': file_type  # Used for syntax highlighting
            })

        except Exception as e:
            logger.error(f"Error reading text file: {e}")
            return HttpResponseServerError("Error reading text file.")

    except Exception as e:
        logger.error(f"Unhandled exception in view_text_file: {e}")
        return HttpResponseServerError("An error occurred while processing your request.")


@login_required
def view_video_file(request, upload_id):
    try:
        # Fetch the uploaded file record for the authenticated user
        upload = get_object_or_404(UserUpload, id=upload_id, user=request.user)

        # Ensure the file has an encryption key
        if not upload.encryption_key:
            logger.error(f"Encryption key not found for file with ID {upload_id}")
            return HttpResponseServerError("Encryption key not found for this file.")

        try:
            # Initialize the Fernet cipher for decryption using the encryption key
            fernet = Fernet(upload.encryption_key.encode('utf-8'))
        except Exception as e:
            logger.error(f"Error initializing Fernet cipher: {e}")
            return HttpResponseServerError("Failed to initialize decryption.")

        # Check if the file exists on the server
        file_path = upload.file.path
        if not os.path.exists(file_path):
            logger.error(f"File not found: {file_path}")
            return HttpResponseNotFound("File not found")

        # Get MIME type of the file
        mime_type, _ = mimetypes.guess_type(file_path)
        supported_video_formats = ['video/mp4', 'video/mov', 'video/quicktime', 'video/mpeg']
        if mime_type and mime_type.lower() not in [fmt.lower() for fmt in supported_video_formats]:
            logger.error(f"Unsupported video file type: {mime_type}")
            return HttpResponseServerError("Unsupported video file type.")

        # Read and decrypt the video file
        try:
            with open(file_path, 'rb') as encrypted_file:
                encrypted_file_data = encrypted_file.read()

            decrypted_file_data = fernet.decrypt(encrypted_file_data)

            # Serve the decrypted file directly as a response
            response = HttpResponse(decrypted_file_data, content_type=mime_type)
            response['Content-Disposition'] = f'inline; filename="{upload.file_name}"'
            return response

        except Exception as e:
            logger.error(f"Error reading or decrypting video file: {e}")
            return HttpResponseServerError("Error reading or decrypting video file.")

    except Exception as e:
        logger.error(f"Unhandled exception in view_video_file: {e}")
        return HttpResponseServerError("An error occurred while processing your request.")
    

@login_required
def upload_folder(request):
    if request.method == 'POST':
        form = UserUploadForm(request.POST, request.FILES)
        folder_files = request.FILES.getlist('folder_files')  # Expecting multiple files
        folder_name = request.POST.get('folder_name', 'uploaded_folder')

        if folder_files and form.is_valid():
            form.instance.user = request.user

            # Step 1: Create a zip archive from the files
            temp_zip_path = f'/tmp/{folder_name}.zip'
            with zipfile.ZipFile(temp_zip_path, 'w') as zipf:
                for file in folder_files:
                    zipf.writestr(file.name, file.read())

            # Step 2: Generate an encryption key and encrypt the archive
            encryption_key = Fernet.generate_key()
            fernet = Fernet(encryption_key)

            with open(temp_zip_path, 'rb') as zip_file:
                encrypted_data = fernet.encrypt(zip_file.read())

            # Step 3: Save the encrypted archive to the model
            form.instance.file = ContentFile(encrypted_data, f'{folder_name}.zip')
            form.instance.file_name = f'{folder_name}.zip'
            form.instance.encryption_key = encryption_key.decode('utf-8')
            form.instance.is_folder = True

            try:
                form.save()
            except ValueError:
                messages.error(request, "Storage full. Please buy more storage.")
                return redirect('/buy_storage')

            # Step 4: Clean up temporary files
            os.remove(temp_zip_path)

            # Redirect after successful upload
            if CustomGroupAdmin.objects.filter(user=request.user).exists():
                return redirect('/subgroup_landing_page')
            else:
                return redirect('/landing_page')
    else:
        form = UserUploadForm()

    return render(request, 'upload_folder.html', {'form': form})



@login_required
def view_folder(request, upload_id):
    try:
        # Fetch the folder record
        upload = get_object_or_404(UserUpload, id=upload_id, user=request.user, is_folder=True)

        if not upload.encryption_key:
            logger.error(f"Encryption key missing for upload ID {upload_id}")
            return HttpResponseServerError("Encryption key missing.")

        # Decrypt the folder archive
        try:
            fernet = Fernet(upload.encryption_key.encode('utf-8'))
            encrypted_data = upload.file.read()
            decrypted_data = fernet.decrypt(encrypted_data)
        except Exception as e:
            logger.error(f"Decryption failed for upload ID {upload_id}: {e}")
            return HttpResponseServerError("Decryption failed.")

        # Extract the archive to a temporary location
        temp_extract_path = f'/tmp/{upload.file_name.replace(".zip", "")}'
        os.makedirs(temp_extract_path, exist_ok=True)

        temp_zip_path = f'{temp_extract_path}.zip'
        with open(temp_zip_path, 'wb') as temp_zip:
            temp_zip.write(decrypted_data)

        with zipfile.ZipFile(temp_zip_path, 'r') as zipf:
            zipf.extractall(temp_extract_path)

        os.remove(temp_zip_path)

        # List extracted files and provide links to view/download
        file_links = []
        for root, _, files in os.walk(temp_extract_path):
            for file in files:
                file_path = os.path.join(root, file)
                file_links.append(file_path.replace('/tmp/', '/tmp/'))

        return render(request, 'view_folder.html', {'files': file_links, 'upload': upload})

    except Exception as e:
        logger.error(f"Unhandled exception in view_folder: {e}")
        return HttpResponseServerError("An error occurred while processing the folder.")
    finally:
        # Clean up extracted files after serving the response
        if 'temp_extract_path' in locals():
            shutil.rmtree(temp_extract_path, ignore_errors=True)




@login_required
def delete_upload(request, upload_id):
    upload = get_object_or_404(UserUpload, id=upload_id)
    if request.method == 'POST':
        form = DeleteUploadForm(request.POST)
        if form.is_valid():
            upload.delete_file()  # Ensure the file is deleted from the filesystem
            upload.delete()  # Delete the database record
            return redirect('/landing_page')
    else:
        form = DeleteUploadForm(initial={'upload_id': upload_id})
    return render(request, 'delete_upload.html', {'form': form})


def upload_file(request):
    if request.method == 'POST':
        uploaded_file = request.FILES['file']
        upload = UserUpload.objects.create(file=uploaded_file)
        return redirect('/upload_success')
    return render(request, 'upload_form.html',{'upload':upload})

def registration_form_view(request):
    if request.method == 'POST':
        form = RegistrationFormForm(request.POST, request.FILES)
        if form.is_valid():
            form.instance.user = request.user
            registration = form.save()
            subgroup_id = request.POST.get('subgroup')
            if subgroup_id:
                try:
                    subgroup = CustomGroup.objects.get(pk=subgroup_id)
                    registration.subgroup = subgroup
                    registration.save()
                except CustomGroup.DoesNotExist:
                    pass
            return redirect('/landing_page')
    else:
        form = RegistrationFormForm()
    return render(request, 'RegistrationForm.html', {'form': form})

@staff_member_required
def pending_associations(request):
    pending_users = UserAssociation.objects.filter(is_approved=False)
    return render(request, 'pending_associations.html', {'pending_users': pending_users})


@staff_member_required
def approve_association(request, association_id):
    association = UserAssociation.objects.get(pk=association_id)
    association.is_approved = True
    association.save()
    return redirect('pending_associations')


@login_required
def send_file(request):
    if request.method == 'POST':
        file = request.FILES.get('file')
        if file:
            File.objects.create(user=request.user, file=file)
            return redirect('/landing_page')
    return render(request, 'send_file.html')


@login_required
def create_card(request):
    if request.method == 'POST':
        form = CardForm(request.POST, request.FILES)
        if form.is_valid():
            try:
                request.user.card
                form.add_error(None, "A card already exists for this user.")
            except User.card.RelatedObjectDoesNotExist:
                new_card = form.save(commit=False)
                new_card.user = request.user
                new_card.save()
                return redirect('/landing_page')
    else:
        form = CardForm()
    return render(request, 'create_card.html', {'card_form': form})


@login_required
def create_kyc(request):
    if request.method == 'POST':
        kyc_form = KYCForm(request.POST, request.FILES)
        if kyc_form.is_valid():
            user = request.user if request.user.is_authenticated else None
            uploaded_image = kyc_form.cleaned_data['uploaded_image']
            kyc_video = kyc_form.cleaned_data['kyc_video']
            phone_number = kyc_form.cleaned_data['phone_number']
            kyc = KYC(user=user, uploaded_image=uploaded_image, kyc_video=kyc_video, phone_number=phone_number)
            kyc.save()
            account_sid = 'AC87e2f3c11c67a2b0a6dd4f644d5900a7'
            auth_token = 'a76b9c54d676468919c4fa5a12785d3d'
            client = Client(account_sid, auth_token)
            try:
                message = client.messages.create(
                    body='KYC verification: Your photo and video have been received. Please wait for verification.',
                    from_='7565830698',
                    to=phone_number
                )
                messages.success(request, 'KYC verification message sent successfully.')
            except Exception as e:
                messages.error(request, f'Failed to send KYC verification message: {e}')
            return redirect('/landing_page')
        else:
            messages.error(request, 'Failed to submit KYC data. Please check the form entries.')
            return redirect('/landing_page')
    else:
        kyc_form = KYCForm()
    return render(request, 'create_kyc.html', {'kyc_form': kyc_form})

@login_required
def service1(request):
    return render(request, 'service1.html')

@login_required
def service2(request):
    return render(request, 'service2.html')

@login_required
def service3(request):
    return render(request, 'service3.html')



# def signup(request):
#     if request.method == 'POST':
#         form = CustomSignupForm(request.POST)
#         if form.is_valid():
#             # Extract user data from the form
#             username = form.cleaned_data.get('username')
#             email = form.cleaned_data.get('email')
#             password = form.cleaned_data.get('password1')  # Assuming you have a password1 field
#             first_name = form.cleaned_data.get('first_name')
#             last_name = form.cleaned_data.get('last_name')

#             # Generate a temporary token for email confirmation
#             confirmation_token = str(uuid.uuid4())

#             try:
#                 # Check if TemporaryUser with the given email exists
#                 temporary_user, created = TemporaryUser.objects.update_or_create(
#                     email=email,
#                     defaults={
#                         'username': username,
#                         'password': password,  # Should be hashed
#                         'first_name': first_name,
#                         'last_name': last_name,
#                         'token': confirmation_token,
#                         'token_expires': timezone.now() + timedelta(hours=24),
#                     }
#                 )

#                 # Send confirmation email with the token
#                 confirmation_url = request.build_absolute_uri(f'/confirm-email/{confirmation_token}/')
#                 send_mail(
#                     'Confirm your email',
#                     f'Please confirm your email by clicking the link: {confirmation_url}',
#                     'no-reply@socyfie.com',  # This should match or be consistent with DEFAULT_FROM_EMAIL in settings
#                     [email],
#                     fail_silently=False,
#                 )

#                 messages.success(request, "A confirmation email has been sent. Please check your inbox.")
#                 return redirect('/')

#             except Exception as e:
#                 messages.error(request, f"Failed to initiate account creation: {e}")
#         else:
#             messages.error(request, "Failed to create account. Please check the form entries.")
#     else:
#         form = CustomSignupForm()

#     return render(request, 'signup.html', {'form': form})


# from django.contrib.auth.hashers import make_password

# def confirm_email(request, token):
#     # Get the temporary user by token
#     temporary_user = get_object_or_404(TemporaryUser, token=token)

#     # Check if token is still valid
#     if not temporary_user.is_token_valid():
#         messages.error(request, "The confirmation link has expired. Please try signing up again.")
#         return redirect('/signup')

#     try:
#         # Create user in Firebase with the stored credentials
#         # If the password is stored hashed, you may need to handle Firebase accordingly
#         firebase_user = auth.create_user_with_email_and_password(
#             temporary_user.email,
#             temporary_user.password  # If stored hashed, adjust this for Firebase requirements
#         )
#         auth.send_email_verification(firebase_user['idToken'])  # Optional Firebase email verification

#         # Now create the user in Django
#         user_model = User(
#             first_name=temporary_user.first_name,
#             last_name=temporary_user.last_name,
#             email=temporary_user.email,
#             username=temporary_user.username,
#         )
#         # If the password was hashed before, apply the same logic here
#         user_model.set_password(temporary_user.password)  # Ensure hashing for Django user password
#         user_model.save()

#         # Log the user in Django (Optional: Consider waiting until Firebase email verification is complete)
#         login(request, user_model)

#         # Remove the temporary user record after successful creation
#         temporary_user.delete()

#         messages.success(request, "Your account has been successfully verified and created.")
#         return redirect('/')

#     except Exception as e:
#         messages.error(request, f"Account creation failed: {e}")
#         return redirect('/signup')

# def check_firebase_email_verification(request, user_id):
#     try:
#         # Get the user from Firebase by their Django user_id or Firebase UID
#         user = User.objects.get(pk=user_id)
        
#         # Re-authenticate or refresh the Firebase user details
#         firebase_user = auth.get_account_info(user.firebase_uid)  # Fetch Firebase user details
        
#         # Check if the email is verified
#         if firebase_user['users'][0]['emailVerified']:
#             # Activate the user in Django
#             user.is_active = True
#             user.save()

#             # Log the user in Django
#             login(request, user)
            
#             messages.success(request, "Your email has been verified. You are now logged in.")
#             return redirect('/')

#         else:
#             messages.error(request, "Your email is not verified yet. Please check your email.")
#             return redirect('/verify-email-pending')

#     except User.DoesNotExist:
#         messages.error(request, "User not found.")
#         return redirect('/signup')

#     except Exception as e:
#         messages.error(request, f"Error checking verification: {e}")
#         return redirect('/signup')
